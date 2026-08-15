"""Optional XLSX and PPTX parsers with structural provenance."""

# pyright: reportMissingImports=false, reportMissingModuleSource=false

from __future__ import annotations

from typing import Any

from app.parsers.base import BaseParser, ParsedDocument


class SpreadsheetParser(BaseParser):
    """Parse XLSX sheets into deterministic row-labelled text."""

    async def parse(self, file_path: str) -> ParsedDocument:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise RuntimeError(
                "XLSX parsing requires the optional groundtruth[office] extra"
            ) from exc

        workbook = load_workbook(file_path, read_only=True, data_only=True)
        rendered: list[str] = []
        sections: list[str] = []
        provenance: list[dict[str, Any]] = []
        row_count = 0
        try:
            for sheet in workbook.worksheets:
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    continue
                headers = [
                    str(value) if value is not None else f"column_{index}"
                    for index, value in enumerate(rows[0], start=1)
                ]
                for sheet_row, values in enumerate(rows[1:], start=2):
                    if not any(value is not None for value in values):
                        continue
                    row_count += 1
                    section = f"Sheet {sheet.title} / Row {sheet_row}"
                    fields = [
                        f"{header}={value if value is not None else ''}"
                        for header, value in zip(headers, values, strict=False)
                    ]
                    rendered.append(f"{section}: " + " | ".join(fields))
                    sections.append(section)
                    provenance.append({"sheet": sheet.title, "row": sheet_row})
        finally:
            workbook.close()

        return ParsedDocument(
            content="\n".join(rendered),
            metadata={
                "file_type": "xlsx",
                "sheet_count": len(workbook.sheetnames),
                "row_count": row_count,
                "chunk_provenance": provenance,
            },
            sections=sections,
        )


class PowerPointParser(BaseParser):
    """Parse PPTX text by slide while preserving slide provenance."""

    async def parse(self, file_path: str) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError(
                "PPTX parsing requires the optional groundtruth[office] extra"
            ) from exc

        presentation = Presentation(file_path)
        rendered: list[str] = []
        sections: list[str] = []
        provenance: list[dict[str, int]] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts = [
                str(shape.text).strip()
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
                and str(getattr(shape, "text", "")).strip()
            ]
            if not texts:
                continue
            section = f"Slide {slide_index}"
            rendered.append(f"{section}: " + "\n".join(texts))
            sections.append(section)
            provenance.append({"slide": slide_index})

        return ParsedDocument(
            content="\n\n".join(rendered),
            metadata={
                "file_type": "pptx",
                "slide_count": len(presentation.slides),
                "chunk_provenance": provenance,
            },
            sections=sections,
        )
